"""
SAM.gov Contract RAG Query Tool
================================
Interactive Q&A grounded in the actual solicitation documents and past
performance decks. Run this after extractor_multi.py to ask follow-up
questions, map past performance to requirements, or draft section responses.

Two persistent ChromaDB collections:
  solicitation_{name}  — chunks from the current RFP package (per opportunity)
  past_performance     — chunks from PP decks (shared across all opportunities)

Usage:
    python interactive_query.py <solicitation_folder> <pp_folder>
    python interactive_query.py <solicitation_folder> <pp_folder> --rebuild
    python interactive_query.py <solicitation_folder> <pp_folder> --rebuild-pp

Example:
    python interactive_query.py "Sample Contracts/US_Secret Service" "Past Performance"

Example questions:
    > What are the key deliverables in the PWS?
    > Which of our past performances best maps to the data analytics requirement?
    > Draft a response to the key personnel requirement using our Army past performance.

Author: Harry Cotton
"""

import anthropic
import chromadb
import sys
import os
import glob
import re


CHROMA_DIR = "./chroma_db"
PP_COLLECTION_NAME = "past_performance"
N_RESULTS = 4        # chunks retrieved per collection per query
CHUNK_SIZE = 3000    # max chars per chunk
CHUNK_OVERLAP = 300  # overlap between chunks to preserve context


# ============================================================
# STEP 1: FILE READERS
# (mirrors extractor_multi.py — kept separate so scripts stay independent)
# ============================================================

OUTPUT_SUFFIXES = (
    "_extraction.txt", "_extracted.txt",
    "_full_extraction.json", "_raw_extraction.txt",
    "_proposal_draft.md", "_past_performance_draft.md",
)


def read_pdf(file_path: str) -> str:
    try:
        import pdfplumber
    except ImportError:
        print("Error: pdfplumber required. Run: pip install pdfplumber")
        sys.exit(1)
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            if page_text:
                text += f"[Page {i + 1}]\n{page_text}\n\n"
    return text


def read_docx(file_path: str) -> str:
    try:
        from docx import Document
    except ImportError:
        print("Error: python-docx required. Run: pip install python-docx")
        sys.exit(1)
    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    return "\n".join(paragraphs)


def read_excel(file_path: str) -> str:
    try:
        import openpyxl
    except ImportError:
        print("Error: openpyxl required. Run: pip install openpyxl")
        sys.exit(1)
    wb = openpyxl.load_workbook(file_path, data_only=True)
    text = ""
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(cells)
        if rows:
            text += f"[Sheet: {sheet_name}]\n"
            for row in rows:
                text += " | ".join(row) + "\n"
            text += "\n"
    return text


def read_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        print("Error: python-pptx required. Run: pip install python-pptx")
        sys.exit(1)
    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        parts = []
        if slide.shapes.title and slide.shapes.title.text.strip():
            parts.append(f"[Slide {i + 1}]: {slide.shapes.title.text.strip()}")
        else:
            parts.append(f"[Slide {i + 1}]")
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
            elif hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Notes]: {notes}")
        if len(parts) > 1:
            slides_text.append("\n".join(parts))
    return "\n\n".join(slides_text)


READERS = {
    ".pdf":  read_pdf,
    ".docx": read_docx,
    ".xlsx": read_excel,
    ".xls":  read_excel,
    ".pptx": read_pptx,
}


# ============================================================
# STEP 2: CHUNKING
# ============================================================

def chunk_text(text: str, source_file: str, doc_type: str) -> tuple:
    """
    Splits text into overlapping chunks.

    For solicitation PDFs the text already has [Page N] markers from
    pdfplumber, so we split on those first to keep page boundaries clean.
    For PP slides we split on [Slide N] markers. Anything without markers
    (e.g. Excel, Word) falls back to plain character chunking.
    """
    documents, metadatas, ids = [], [], []
    filename = os.path.basename(source_file)
    safe_name = re.sub(r"[^a-zA-Z0-9]", "_", filename)

    # Try splitting on [Page N] or [Slide N] markers
    page_pattern = re.compile(r"(\[(?:Page|Slide) \d+[^\]]*\])")
    parts = page_pattern.split(text)

    if len(parts) > 1:
        # Pair each marker with its following content
        chunks = []
        for i in range(1, len(parts), 2):
            marker = parts[i]
            content = parts[i + 1] if i + 1 < len(parts) else ""
            full = (marker + "\n" + content).strip()
            if full:
                chunks.append((marker, full))
    else:
        # Plain character chunking with overlap
        chunks = []
        start = 0
        idx = 0
        while start < len(text):
            end = start + CHUNK_SIZE
            chunk = text[start:end].strip()
            if chunk:
                chunks.append((f"chunk_{idx}", chunk))
            start = end - CHUNK_OVERLAP
            idx += 1

    for i, (marker, content) in enumerate(chunks):
        # Extract page/slide number from marker if present
        num_match = re.search(r"\d+", marker)
        page_num = int(num_match.group()) if num_match else i + 1

        chunk_id = f"{safe_name}_{doc_type}_{i}"
        documents.append(content)
        metadatas.append({
            "source_file": filename,
            "doc_type": doc_type,
            "page_or_slide": page_num,
        })
        ids.append(chunk_id)

    return documents, metadatas, ids


# ============================================================
# STEP 3: BUILD VECTOR STORES
# ============================================================

def safe_collection_name(folder_path: str) -> str:
    name = os.path.basename(os.path.normpath(folder_path))
    safe = re.sub(r"[^a-zA-Z0-9_-]", "_", name).strip("_")
    return f"solicitation_{safe[:40]}"


def build_solicitation_store(folder_path: str, chroma_client, collection_name: str,
                             force_rebuild: bool = False) -> chromadb.Collection:
    if not force_rebuild:
        try:
            col = chroma_client.get_collection(collection_name)
            print(f"Loaded solicitation store: {col.count()} chunks")
            return col
        except Exception:
            pass

    print(f"Building solicitation vector store...")

    supported_exts = list(READERS.keys())
    all_files = []
    for ext in supported_exts:
        for f in glob.glob(os.path.join(folder_path, f"*{ext}")):
            basename = os.path.basename(f)
            if not any(basename.endswith(s) for s in OUTPUT_SUFFIXES):
                all_files.append(f)
    all_files.sort()

    if not all_files:
        print(f"Error: No supported documents found in {folder_path}")
        sys.exit(1)

    print(f"  Found {len(all_files)} document(s): {', '.join(os.path.basename(f) for f in all_files)}")

    all_docs, all_meta, all_ids = [], [], []
    for file_path in all_files:
        ext = os.path.splitext(file_path)[1].lower()
        reader = READERS[ext]
        print(f"  Chunking: {os.path.basename(file_path)}...")
        text = reader(file_path)
        docs, metas, ids = chunk_text(text, file_path, "solicitation")
        all_docs.extend(docs)
        all_meta.extend(metas)
        all_ids.extend(ids)

    try:
        chroma_client.delete_collection(collection_name)
    except Exception:
        pass

    col = chroma_client.create_collection(collection_name)
    # Add in batches to avoid ChromaDB limits
    batch_size = 100
    for i in range(0, len(all_docs), batch_size):
        col.add(
            documents=all_docs[i:i + batch_size],
            metadatas=all_meta[i:i + batch_size],
            ids=all_ids[i:i + batch_size],
        )

    print(f"  Built solicitation store: {col.count()} chunks\n")
    return col


def build_pp_store(pp_folder: str, chroma_client,
                   force_rebuild: bool = False) -> chromadb.Collection:
    if not force_rebuild:
        try:
            col = chroma_client.get_collection(PP_COLLECTION_NAME)
            print(f"Loaded past performance store: {col.count()} chunks")
            return col
        except Exception:
            pass

    print("Building past performance vector store...")

    pptx_files = sorted(glob.glob(os.path.join(pp_folder, "*.pptx")))
    if not pptx_files:
        print(f"Warning: No .pptx files found in {pp_folder}. PP retrieval unavailable.")
        try:
            chroma_client.delete_collection(PP_COLLECTION_NAME)
        except Exception:
            pass
        col = chroma_client.create_collection(PP_COLLECTION_NAME)
        return col

    print(f"  Found {len(pptx_files)} deck(s): {', '.join(os.path.basename(f) for f in pptx_files)}")

    all_docs, all_meta, all_ids = [], [], []
    for file_path in pptx_files:
        print(f"  Chunking: {os.path.basename(file_path)}...")
        text = read_pptx(file_path)
        docs, metas, ids = chunk_text(text, file_path, "past_performance")
        all_docs.extend(docs)
        all_meta.extend(metas)
        all_ids.extend(ids)

    try:
        chroma_client.delete_collection(PP_COLLECTION_NAME)
    except Exception:
        pass

    col = chroma_client.create_collection(PP_COLLECTION_NAME)
    batch_size = 100
    for i in range(0, len(all_docs), batch_size):
        col.add(
            documents=all_docs[i:i + batch_size],
            metadatas=all_meta[i:i + batch_size],
            ids=all_ids[i:i + batch_size],
        )

    print(f"  Built past performance store: {col.count()} chunks\n")
    return col


# ============================================================
# STEP 4: QUERY AND GENERATE
# ============================================================

def query_and_generate(question: str, sol_col: chromadb.Collection,
                       pp_col: chromadb.Collection) -> str:
    """
    Queries both collections, merges results, and generates a grounded answer.
    """
    context_parts = []

    # Query solicitation collection
    if sol_col.count() > 0:
        sol_results = sol_col.query(query_texts=[question], n_results=min(N_RESULTS, sol_col.count()))
        for doc, meta in zip(sol_results["documents"][0], sol_results["metadatas"][0]):
            context_parts.append(
                f"[SOURCE: {meta['source_file']} — Page/Slide {meta['page_or_slide']}]\n{doc}"
            )

    # Query past performance collection
    if pp_col.count() > 0:
        pp_results = pp_col.query(query_texts=[question], n_results=min(N_RESULTS, pp_col.count()))
        for doc, meta in zip(pp_results["documents"][0], pp_results["metadatas"][0]):
            context_parts.append(
                f"[SOURCE: {meta['source_file']} — Slide {meta['page_or_slide']} (PAST PERFORMANCE)]\n{doc}"
            )

    if not context_parts:
        return "No relevant content found in the vector store."

    context = "\n\n" + ("\n\n" + "=" * 50 + "\n\n").join(context_parts)

    client = anthropic.Anthropic()

    system_prompt = """You are a senior BD analyst and proposal writer supporting a government
contracting firm. You answer questions grounded strictly in the provided source
documents — which include solicitation documents (RFPs, PWS, SOWs) and past
performance decks.

Rules:
1. Base every answer solely on the provided context. Do not draw on outside knowledge.
2. Cite your sources inline: (Source: filename, Page/Slide N).
3. When asked to map past performance to a requirement, explicitly name the
   requirement and explain specifically why that past performance satisfies it.
4. When asked to draft proposal language, write polished, near-finished prose
   in first-person plural ("we / our team") — not a template or outline.
5. If the context does not contain enough information to answer fully, say so
   and indicate what additional documents would help.
6. Never use the words "ensure" or "ensuring"."""

    user_prompt = f"""Question: {question}

Context (from solicitation documents and past performance decks):
{context}

Answer grounded in the provided context, citing sources inline."""

    message = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=2048,
        system=system_prompt,
        messages=[{"role": "user", "content": user_prompt}],
    )

    return message.content[0].text


# ============================================================
# STEP 5: MAIN
# ============================================================

def main():
    if len(sys.argv) < 3:
        print("Usage: python interactive_query.py <solicitation_folder> <pp_folder> [--rebuild] [--rebuild-pp]")
        print()
        print('Example: python interactive_query.py "Sample Contracts/US_Secret Service" "Past Performance"')
        sys.exit(1)

    sol_folder = sys.argv[1]
    pp_folder  = sys.argv[2]
    rebuild_sol = "--rebuild"    in sys.argv
    rebuild_pp  = "--rebuild-pp" in sys.argv

    if not os.path.isdir(sol_folder):
        print(f"Error: Solicitation folder not found: {sol_folder}")
        sys.exit(1)
    if not os.path.isdir(pp_folder):
        print(f"Error: Past performance folder not found: {pp_folder}")
        sys.exit(1)

    print(f"Solicitation: {sol_folder}")
    print(f"Past Performance: {pp_folder}")
    print("-" * 60)

    chroma_client = chromadb.PersistentClient(path=CHROMA_DIR)
    collection_name = safe_collection_name(sol_folder)

    sol_col = build_solicitation_store(sol_folder, chroma_client, collection_name, rebuild_sol)
    pp_col  = build_pp_store(pp_folder, chroma_client, rebuild_pp)

    print()
    print("=" * 60)
    print("CONTRACT Q&A — Interactive Mode")
    print("=" * 60)
    print("Ask questions about the solicitation or past performance.")
    print("Examples:")
    print('  "What are the key deliverables in the PWS?"')
    print('  "Which past performance best maps to the data analytics requirement?"')
    print('  "Draft a response to the key personnel section using our Army PP."')
    print("Type 'quit' to exit.\n")

    while True:
        try:
            question = input("Your question: ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\nGoodbye!")
            break

        if not question:
            continue
        if question.lower() in ("quit", "exit", "q"):
            print("Goodbye!")
            break

        print("\nSearching documents...\n")
        answer = query_and_generate(question, sol_col, pp_col)

        print("-" * 60)
        print(answer)
        print("-" * 60 + "\n")


if __name__ == "__main__":
    main()
