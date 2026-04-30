"""
Proposal Writer
===============
Takes the JSON extraction output from extractor_multi.py and a folder of
past performance PowerPoint decks, then uses Claude to produce a near-finished
full proposal response structured around the proposal outline identified in
the extraction JSON.

Unlike pp_writer.py (which drafts the past performance volume in isolation),
this script produces a complete proposal draft -- each volume written as
near-finished prose, with past performance examples woven in as evidence
supporting the technical and management claims. The extraction JSON is the
blueprint; the PP decks are the proof.

Usage:
    python proposal_writer.py <extraction_json> <past_performance_folder>

Example:
    python proposal_writer.py "Sample Contracts/US_Secret Service/US_Secret Service_full_extraction.json" "Past Performance"

Output:
    A markdown file saved alongside the extraction JSON, structured by
    proposal volume and ready for staff review and light editing.

Required libraries:
    pip install anthropic python-pptx

Author: Harry Cotton
"""

import anthropic
import sys
import os
import json
import glob
from pathlib import Path


# ============================================================
# STEP 1: EXTRACT TEXT FROM POWERPOINT DECKS
# (shared logic with pp_writer.py)
# ============================================================

def extract_text_from_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        print("Error: python-pptx is required. Run: pip install python-pptx")
        sys.exit(1)

    prs = Presentation(file_path)
    slides_text = []

    for i, slide in enumerate(prs.slides):
        slide_content = []

        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_content.append(f"[Slide {i + 1}]: {slide.shapes.title.text.strip()}")
        else:
            slide_content.append(f"[Slide {i + 1}]")

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_content.append(" | ".join(cells))
            elif hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_content.append(f"[Notes]: {notes}")

        if len(slide_content) > 1:
            slides_text.append("\n".join(slide_content))

    return "\n\n".join(slides_text)


def read_past_performance_decks(folder_path: str) -> list:
    pptx_files = sorted(glob.glob(os.path.join(folder_path, "*.pptx")))

    if not pptx_files:
        print(f"Error: No .pptx files found in {folder_path}")
        sys.exit(1)

    decks = []
    print(f"Found {len(pptx_files)} past performance deck(s):")
    for f in pptx_files:
        print(f"  - {os.path.basename(f)}")
        text = extract_text_from_pptx(f)
        decks.append({"filename": os.path.basename(f), "text": text})

    return decks


def format_decks_for_prompt(decks: list) -> str:
    combined = ""
    for deck in decks:
        combined += f"\n{'='*60}\n"
        combined += f"PAST PERFORMANCE DECK: {deck['filename']}\n"
        combined += f"{'='*60}\n\n"
        combined += deck["text"]
        combined += "\n\n"
    return combined


# ============================================================
# STEP 2: LOAD FIRM CONFIG
# ============================================================

def load_firm_config() -> dict:
    config_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "firm_config.json")

    if not os.path.exists(config_path):
        print("Warning: firm_config.json not found. Using placeholder firm details.")
        return {
            "firm_name": "Our Firm",
            "firm_description": "a government contracting firm",
            "firm_capabilities": "consulting, program management, and technology services",
            "certifications": [],
            "naics_codes": []
        }

    with open(config_path, "r", encoding="utf-8") as f:
        config = json.load(f)

    print(f"Firm: {config.get('firm_name', 'Unknown')}")
    return config


# ============================================================
# STEP 3: BUILD FULL SOLICITATION CONTEXT FROM EXTRACTION JSON
# ============================================================

def build_full_solicitation_context(extraction: dict) -> str:
    """
    Pulls everything from the extraction JSON that informs proposal strategy:
    overview, requirements, eval criteria, competitive intelligence, risks,
    and the full proposal outline with section-level guidance.
    """
    overview = extraction.get("opportunity_overview", {})
    details = extraction.get("contract_details", {})
    requirements = extraction.get("requirements_summary", {})
    eval_criteria = extraction.get("evaluation_criteria", {})
    intel = extraction.get("competitive_intelligence", {})
    actions = extraction.get("bd_action_items", {})
    outline = extraction.get("proposal_outline", {})

    context = f"""OPPORTUNITY: {overview.get('title', 'N/A')}
AGENCY: {overview.get('agency', 'N/A')}
SOLICITATION #: {overview.get('solicitation_number', 'N/A')}
RESPONSE DEADLINE: {overview.get('response_deadline', 'N/A')}
CONTRACT TYPE: {details.get('contract_type', 'N/A')}
SET-ASIDE: {details.get('set_aside', 'N/A')}
ESTIMATED VALUE: {details.get('estimated_value', 'N/A')}
PERIOD OF PERFORMANCE: {details.get('period_of_performance', 'N/A')}

SCOPE:
{requirements.get('scope_overview', 'N/A')}

KEY REQUIREMENTS:
{chr(10).join(f'  - {r}' for r in requirements.get('key_requirements', []))}

EVALUATION METHOD: {eval_criteria.get('evaluation_method', 'N/A')}

EVALUATION FACTORS:
{chr(10).join(f'  - {f}' for f in eval_criteria.get('factors', []))}

SUBMISSION REQUIREMENTS:
{eval_criteria.get('submission_requirements', 'N/A')}

COMPETITIVE INTELLIGENCE:
  Incumbent: {intel.get('incumbent', 'N/A')}
  Competition: {intel.get('estimated_competition', 'N/A')}
  Notable Terms: {intel.get('notable_terms', 'N/A')}

RISKS AND FLAGS TO ADDRESS IN THE PROPOSAL:
{chr(10).join(f'  - {r}' for r in actions.get('risks_and_flags', []))}

PURSUIT RECOMMENDATION:
{actions.get('pursuit_recommendation', 'N/A')}

WIN THEMES (thread throughout every section):
{chr(10).join(f'  - {t}' for t in outline.get('win_themes', []))}

PROPOSAL OUTLINE — WRITE EACH SECTION BELOW:
"""

    for i, section in enumerate(outline.get("sections", []), 1):
        context += f"""
SECTION {i}: {section.get('title', 'Untitled')}
  Requirements this section must address:
{chr(10).join(f'    - {r}' for r in section.get('requirements_addressed', []))}
  Guidance: {section.get('guidance', 'N/A')}
"""

    return context


# ============================================================
# STEP 4: DRAFT THE FULL PROPOSAL
# ============================================================

def write_full_proposal(extraction: dict, decks: list, firm_config: dict) -> str:
    """
    Sends the full solicitation context and PP decks to Claude.
    Claude writes a complete, near-finished proposal response structured
    around the proposal outline from the extraction JSON, with PP examples
    woven in as evidence throughout.
    """
    solicitation_context = build_full_solicitation_context(extraction)
    decks_text = format_decks_for_prompt(decks)
    firm_name = firm_config.get("firm_name", "Our Firm")
    firm_description = firm_config.get("firm_description", "")
    firm_capabilities = firm_config.get("firm_capabilities", "")
    certs = ", ".join(firm_config.get("certifications", []))

    client = anthropic.Anthropic()

    system_prompt = f"""You are a senior proposal writer at {firm_name} with 15 years of experience
writing winning federal proposals. You write near-finished, submission-ready prose
that staff can lightly tune -- not templates, not outlines, not placeholders.

ABOUT {firm_name.upper()}:
{firm_name} is {firm_description}
Core capabilities: {firm_capabilities}
Certifications: {certs if certs else 'See firm profile'}
{firm_name} has the depth and breadth of a Big Four consulting firm, with a proven
track record delivering complex, mission-critical programs for federal agencies.

YOUR TASK:
Write a complete proposal response to the solicitation described below. Structure
the response exactly according to the PROPOSAL OUTLINE provided -- one section per
proposal volume. Each section must:

1. Be written as polished, near-finished proposal prose from {firm_name}'s perspective
   using "we," "our team," and "{firm_name}" throughout
2. Directly address every requirement mapped to that section in the proposal outline
3. Weave in specific past performance examples from the provided decks as proof points
   -- cite them inline where they demonstrate a claimed capability, not as a separate
   appendix or afterthought
4. Thread the win themes throughout -- every section should reinforce the same
   overarching differentiators
5. Proactively address any relevant risks or flags identified in the solicitation
   analysis where doing so strengthens the proposal

PAST PERFORMANCE INTEGRATION:
Do not write a separate past performance section as a list of narratives. Instead:
- In the Technical section, cite PP examples that prove you have done this before
- In the Management section, cite PP examples that demonstrate your QCP, key personnel,
  and phase-in approach
- In the dedicated Past Performance section (if one exists in the outline), write
  full structured narratives in standard govcon format
- Every capability claim should be backed by at least one reference to a specific
  past project, metric, or outcome from the PP decks

BANNED WORDS -- never use these anywhere in the output:
- "ensure" or "ensuring"

FORMAT:
- Use clear markdown headers matching the proposal volume titles from the outline
- Write in third person / first person plural ("we / our / {firm_name}")
- Use professional proposal language -- active voice, specific, confident
- Flag missing data as [TO BE CONFIRMED] so staff know what to verify
- After all proposal sections, include a brief "Proposal Team Checklist" of the
  top 5 items staff must verify or complete before submission"""

    user_prompt = f"""Write a complete, near-finished proposal response for {firm_name} using the
solicitation context and past performance decks below.

Structure the response exactly according to the PROPOSAL OUTLINE sections listed in
the solicitation context. For each section, write polished proposal prose that
addresses the mapped requirements, weaves in relevant PP examples as proof, and
threads the win themes throughout.

This should read like a draft that needs light editing -- not a template.

--- SOLICITATION CONTEXT AND PROPOSAL OUTLINE ---
{solicitation_context}

--- PAST PERFORMANCE DECKS ---
{decks_text}

Write the complete proposal response now, section by section, followed by the
Proposal Team Checklist."""

    message = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=16000,
        system=system_prompt,
        messages=[{"role": "user", "content": user_prompt}]
    )

    return message.content[0].text


# ============================================================
# STEP 5: MAIN
# ============================================================

def main():
    if len(sys.argv) < 3:
        print("Usage: python proposal_writer.py <extraction_json> <past_performance_folder>")
        print()
        print("Example:")
        print('  python proposal_writer.py "Sample Contracts/US_Secret Service/US_Secret Service_full_extraction.json" "Past Performance"')
        sys.exit(1)

    json_path = sys.argv[1]
    pp_folder = sys.argv[2]

    if not os.path.exists(json_path):
        print(f"Error: Extraction JSON not found: {json_path}")
        sys.exit(1)

    if not os.path.isdir(pp_folder):
        print(f"Error: Past performance folder not found: {pp_folder}")
        sys.exit(1)

    # Load firm config and extraction JSON
    firm_config = load_firm_config()

    with open(json_path, "r", encoding="utf-8") as f:
        extraction = json.load(f)

    opportunity_title = extraction.get("opportunity_overview", {}).get("title", "Unknown Opportunity")

    print(f"Opportunity:  {opportunity_title}")
    print(f"Extraction:   {json_path}")
    print(f"PP folder:    {pp_folder}")
    print("-" * 60)

    # Read PP decks
    decks = read_past_performance_decks(pp_folder)
    print()

    # Draft the full proposal
    print("Drafting full proposal response...\n")
    draft = write_full_proposal(extraction, decks, firm_config)

    # Save as markdown alongside the extraction JSON
    json_dir = os.path.dirname(os.path.abspath(json_path))
    stem = Path(json_path).stem.replace("_full_extraction", "")
    output_path = os.path.join(json_dir, f"{stem}_proposal_draft.md")

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(f"# Proposal Draft\n")
        f.write(f"**Opportunity:** {opportunity_title}  \n")
        f.write(f"**Firm:** {firm_config.get('firm_name', 'Unknown')}  \n")
        f.write(f"**Source:** {os.path.basename(json_path)}  \n\n")
        f.write("---\n\n")
        f.write(draft)

    print("=" * 60)
    print(draft)
    print("=" * 60)
    print(f"\nDraft saved to: {output_path}")


if __name__ == "__main__":
    main()
