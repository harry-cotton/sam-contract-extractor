"""
Past Performance Writer
=======================
Takes the JSON extraction output from extractor_multi.py (or extractor.py)
and a folder of past performance PowerPoint decks, then uses Claude to:

1. Review the solicitation requirements and past performance evaluation criteria
2. Select the 3-5 most relevant past performance examples from your decks
3. Draft each in standard govcon past performance narrative format, with
   explicit relevance statements tied to the evaluation criteria

Usage:
    python pp_writer.py <extraction_json> <past_performance_folder>

Example:
    python pp_writer.py "Sample Contracts/US_Secret Service/US_Secret Service_full_extraction.json" "Past Performance"

Output:
    A markdown file saved alongside the extraction JSON, ready to edit and
    drop into your proposal.

Required libraries:
    pip install anthropic python-pptx

Author: Harry Cotton
"""

import anthropic
import sys
import os
import json
import glob
from pathlib import Path


# ============================================================
# STEP 1: EXTRACT TEXT FROM POWERPOINT DECKS
# ============================================================

def extract_text_from_pptx(file_path: str) -> str:
    """
    Pulls all readable text from a PowerPoint file: slide titles, text boxes,
    tables, and speaker notes. Labels each slide so Claude can follow the
    structure of the deck.
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("Error: python-pptx is required. Run: pip install python-pptx")
        sys.exit(1)

    prs = Presentation(file_path)
    slides_text = []

    for i, slide in enumerate(prs.slides):
        slide_content = []

        # Slide title
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_content.append(f"[Slide {i + 1}]: {slide.shapes.title.text.strip()}")
        else:
            slide_content.append(f"[Slide {i + 1}]")

        # All other shapes
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_content.append(" | ".join(cells))
            elif hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())

        # Speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_content.append(f"[Notes]: {notes}")

        if len(slide_content) > 1:
            slides_text.append("\n".join(slide_content))

    return "\n\n".join(slides_text)


def read_past_performance_decks(folder_path: str) -> list:
    """
    Reads all .pptx files in the past performance folder.
    Returns a list of dicts: {filename, text}.
    """
    pptx_files = sorted(glob.glob(os.path.join(folder_path, "*.pptx")))

    if not pptx_files:
        print(f"Error: No .pptx files found in {folder_path}")
        sys.exit(1)

    decks = []
    print(f"Found {len(pptx_files)} past performance deck(s):")
    for f in pptx_files:
        print(f"  - {os.path.basename(f)}")
        text = extract_text_from_pptx(f)
        decks.append({"filename": os.path.basename(f), "text": text})

    return decks


def format_decks_for_prompt(decks: list) -> str:
    combined = ""
    for deck in decks:
        combined += f"\n{'='*60}\n"
        combined += f"PAST PERFORMANCE DECK: {deck['filename']}\n"
        combined += f"{'='*60}\n\n"
        combined += deck["text"]
        combined += "\n\n"
    return combined


# ============================================================
# STEP 2: BUILD SOLICITATION CONTEXT FROM EXTRACTION JSON
# ============================================================

def build_solicitation_context(extraction: dict) -> tuple:
    """
    Pulls the fields from the extraction JSON that matter most for
    past performance selection: scope, key requirements, eval criteria,
    and the proposal outline guidance for the PP volume.
    """
    opportunity = extraction.get("opportunity_overview", {})
    requirements = extraction.get("requirements_summary", {})
    eval_criteria = extraction.get("evaluation_criteria", {})
    outline = extraction.get("proposal_outline", {})

    # Find the past performance section in the proposal outline
    pp_outline = next(
        (s for s in outline.get("sections", [])
         if "past performance" in s.get("title", "").lower()),
        {}
    )

    context = f"""OPPORTUNITY: {opportunity.get('title', 'N/A')}
AGENCY: {opportunity.get('agency', 'N/A')}
SOLICITATION #: {opportunity.get('solicitation_number', 'N/A')}
SCOPE: {requirements.get('scope_overview', 'N/A')}

KEY REQUIREMENTS:
{chr(10).join(f'  - {r}' for r in requirements.get('key_requirements', []))}

PAST PERFORMANCE EVALUATION CRITERIA:
{chr(10).join(f'  - {f}' for f in eval_criteria.get('factors', [])
              if 'past performance' in f.lower()) or '  Not specified'}

PROPOSAL OUTLINE GUIDANCE FOR PAST PERFORMANCE VOLUME:
  Requirements to address: {', '.join(pp_outline.get('requirements_addressed', ['Not specified']))}
  Guidance: {pp_outline.get('guidance', 'Not specified')}"""

    return context, opportunity.get("title", "Unknown Opportunity")


# ============================================================
# STEP 3: LOAD FIRM CONFIG
# ============================================================

def load_firm_config() -> dict:
    """
    Loads firm_config.json from the same directory as this script.
    If not found, returns placeholder values with a warning.
    """
    config_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "firm_config.json")

    if not os.path.exists(config_path):
        print("Warning: firm_config.json not found. Using placeholder firm details.")
        print("         Create firm_config.json in the same folder as pp_writer.py to personalise output.\n")
        return {
            "firm_name": "Our Firm",
            "firm_description": "a government contracting firm",
            "firm_capabilities": "consulting, program management, and technology services",
            "certifications": [],
            "naics_codes": []
        }

    with open(config_path, "r", encoding="utf-8") as f:
        config = json.load(f)

    print(f"Firm: {config.get('firm_name', 'Unknown')}")
    return config


def build_firm_context(config: dict) -> str:
    """Formats firm config into a concise context string for the prompt."""
    name = config.get("firm_name", "Our Firm")
    description = config.get("firm_description", "")
    capabilities = config.get("firm_capabilities", "")
    certs = config.get("certifications", [])
    naics = config.get("naics_codes", [])

    context = f"FIRM NAME: {name}\n"
    if description:
        context += f"FIRM DESCRIPTION: {name} is {description}.\n"
    if capabilities:
        context += f"CORE CAPABILITIES: {capabilities}\n"
    if certs:
        context += f"CERTIFICATIONS: {', '.join(certs)}\n"
    if naics:
        context += f"NAICS CODES: {', '.join(naics)}\n"

    return context


# ============================================================
# STEP 4: DRAFT THE PAST PERFORMANCE VOLUME
# ============================================================

def write_past_performance(extraction: dict, decks: list, firm_config: dict) -> str:
    """
    Sends the solicitation context, firm details, and PP decks to Claude.
    Claude selects the most relevant examples and drafts near-finished
    proposal prose written from the firm's perspective.
    """
    solicitation_context, _ = build_solicitation_context(extraction)
    firm_context = build_firm_context(firm_config)
    decks_text = format_decks_for_prompt(decks)
    firm_name = firm_config.get("firm_name", "Our Firm")

    client = anthropic.Anthropic()

    system_prompt = f"""You are a senior proposal writer at {firm_name} with 15 years of experience
writing winning past performance volumes for federal government contracts.

ABOUT {firm_name.upper()}:
{firm_context}
{firm_name} has the same depth of expertise and capability as the leading Big 4 consulting
firms, with a proven track record of delivering complex, mission-critical programs for
federal agencies.

YOUR ROLE:
You are writing on behalf of {firm_name}. All narratives must be written in first-person
plural ("we," "our team," "{firm_name}") and position the firm as the clear, capable choice
for this opportunity.

WHAT YOU ARE PRODUCING:
A near-finished past performance proposal volume. This is not a template or an outline --
it is polished, submission-ready prose that staff can review and lightly tune before
sending. Every sentence should read as if it came from {firm_name}'s own proposal team.

HOW EVALUATORS SCORE PAST PERFORMANCE:
- Recency: within the last 3-5 years
- Relevance: similar scope, size, and complexity -- larger and more similar is better
- Quality: measurable outcomes, client satisfaction, on-time/on-budget delivery
- Specificity: generic statements do not score -- every relevance claim must name the
  evaluation criterion it satisfies and explain exactly why this project demonstrates it

STANDARD FORMAT FOR EACH ENTRY:
  Project Title:
  Client / Agency:
  Contract Number:        [TO BE CONFIRMED if not in source material]
  Contract Value:         [TO BE CONFIRMED if not in source material]
  Period of Performance:  [TO BE CONFIRMED if not in source material]
  Point of Contact:       [TO BE CONFIRMED if not in source material]

  Project Overview:
  [3-4 sentences of polished prose describing what {firm_name} delivered --
  written as if {firm_name} is presenting this to an evaluator]

  Relevance to This Opportunity:
  [A full paragraph explicitly naming the evaluation criteria this project satisfies
  and explaining exactly how {firm_name}'s experience here directly demonstrates
  the capability required. Do not use vague language.]

  Key Accomplishments:
  [Quantified bullet points. Use real numbers from the source material wherever
  possible. Flag as [TO BE CONFIRMED] if the deck does not include a specific metric.]

BANNED WORDS — NEVER use these words anywhere in the output, including in
headings, narratives, relevance statements, and bullet points:
- "ensure" (compliance flag)
- "ensuring" (compliance flag)

If you would naturally write one of these words, rewrite the sentence to
avoid it entirely. Do not substitute with a close synonym that carries the
same meaning in a hollow way.

CLOSING SECTION:
After the narratives, include:
1. A short "Gaps and Recommendations" section identifying any evaluation criteria
   not well-covered by the selected examples, with specific suggestions for
   how {firm_name} should address the gap (additional references, teaming partners, etc.)
2. A "Proposal Team Notes" section with 3-5 specific items that staff should
   verify, update, or confirm before submission"""

    user_prompt = f"""Using the solicitation context and past performance decks below, select the
3-5 most relevant past performance examples and write the complete, near-finished
past performance proposal volume for {firm_name}.

Write as {firm_name}'s proposal team. Use "{firm_name}" and "we/our" throughout.
Every narrative should read as polished, submission-ready prose -- not a template.

--- SOLICITATION CONTEXT ---
{solicitation_context}

--- PAST PERFORMANCE DECKS ---
{decks_text}

Produce complete narratives for the 3-5 most relevant examples, followed by
the Gaps and Recommendations and Proposal Team Notes sections."""

    message = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=8192,
        system=system_prompt,
        messages=[{"role": "user", "content": user_prompt}]
    )

    return message.content[0].text


# ============================================================
# STEP 4: MAIN
# ============================================================

def main():
    if len(sys.argv) < 3:
        print("Usage: python pp_writer.py <extraction_json> <past_performance_folder>")
        print()
        print("Example:")
        print('  python pp_writer.py "Sample Contracts/US_Secret Service/US_Secret Service_full_extraction.json" "Past Performance"')
        print()
        print("The past performance folder should contain .pptx files, one per project.")
        sys.exit(1)

    json_path = sys.argv[1]
    pp_folder = sys.argv[2]

    if not os.path.exists(json_path):
        print(f"Error: Extraction JSON not found: {json_path}")
        sys.exit(1)

    if not os.path.isdir(pp_folder):
        print(f"Error: Past performance folder not found: {pp_folder}")
        sys.exit(1)

    # Load firm config and extraction JSON
    firm_config = load_firm_config()

    with open(json_path, "r", encoding="utf-8") as f:
        extraction = json.load(f)

    _, opportunity_title = build_solicitation_context(extraction)

    print(f"Opportunity:  {opportunity_title}")
    print(f"Extraction:   {json_path}")
    print(f"PP folder:    {pp_folder}")
    print("-" * 60)

    # Read PP decks
    decks = read_past_performance_decks(pp_folder)
    print()

    # Draft the past performance volume
    print("Drafting past performance narratives...\n")
    draft = write_past_performance(extraction, decks, firm_config)

    # Save as markdown alongside the extraction JSON
    json_dir = os.path.dirname(os.path.abspath(json_path))
    stem = Path(json_path).stem.replace("_full_extraction", "")
    output_path = os.path.join(json_dir, f"{stem}_past_performance_draft.md")

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(f"# Past Performance Draft\n")
        f.write(f"**Opportunity:** {opportunity_title}  \n")
        f.write(f"**Source:** {os.path.basename(json_path)}  \n\n")
        f.write("---\n\n")
        f.write(draft)

    print("=" * 60)
    print(draft)
    print("=" * 60)
    print(f"\nDraft saved to: {output_path}")


if __name__ == "__main__":
    main()
