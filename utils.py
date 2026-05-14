"""
Shared utilities for the SAM.gov BD pipeline.

Functions previously duplicated across pp_writer.py and proposal_writer.py are
centralized here. New helpers (banned-words check, word-budget calculator,
horizontal-rule rendering) live here too so that future fixes land in one
place instead of two.

Modules:
  PowerPoint extraction    — extract_text_from_pptx, read_past_performance_decks
  Prompt formatting        — format_decks_for_prompt
  Markdown -> docx render  — markdown_to_docx, _add_inline_bold
  Compliance checks        — check_banned_words
  Sizing                   — compute_word_budget

Author: Harry Cotton
"""

import os
import re
import sys
import glob


# ============================================================
# Compliance / banned-words check
# ============================================================

# Words that must not appear in generated proposal output.
# Source: every federal-proposal compliance checklist ever. "Ensure" reads as
# legalese in evaluator scoring and is consistently called out.
BANNED_WORDS = ("ensure", "ensuring")

# Anchored to word boundaries so "ensuring" is caught but "insurance" is not.
_BANNED_RE = re.compile(r"\b(?:" + "|".join(BANNED_WORDS) + r")\b", re.IGNORECASE)


def check_banned_words(text: str, label: str = "draft") -> list:
    """
    Scan generated text for banned-word violations. Returns a list of dicts:
        [{"word": "ensuring", "line": 42, "context": "...full line text..."}]

    Does NOT modify the input. Caller decides what to do with the findings
    (re-prompt, warn, log, fail the build, etc.).
    """
    findings = []
    for i, line in enumerate(text.splitlines(), start=1):
        for match in _BANNED_RE.finditer(line):
            findings.append({
                "word": match.group(0).lower(),
                "line": i,
                "context": line.strip()[:180],
            })
    return findings


def warn_on_banned_words(findings: list, label: str = "draft") -> None:
    """Print a clear human-readable warning for any banned-word findings."""
    if not findings:
        return
    print()
    print("=" * 60)
    print(f"WARNING - BANNED WORDS DETECTED in {label}")
    print(f"  {len(findings)} violation(s) found")
    print("=" * 60)
    for f in findings:
        print(f"  Line {f['line']}: '{f['word']}'")
        print(f"    Context: {f['context']}")
    print("=" * 60)
    print("Action required: review the lines above and rewrite to avoid the")
    print("banned word. The draft has been saved but should be revised before")
    print("submission.")
    print("=" * 60)


# ============================================================
# Sizing / page-budget calculator
# ============================================================

# Rough words-per-page in 10pt Times New Roman with 1" margins and single
# spacing. Empirically validated against rendered drafts.
WORDS_PER_PAGE_TNR_10PT_SINGLE = 600


def compute_word_budget(page_limit_total) -> int:
    """
    Convert a Section L page-limit value (number, or string like '20 pages',
    or '30 pages excluding cover/TOC') into a target word count.

    Returns 0 if the page limit cannot be parsed (treat as no budget).
    """
    if page_limit_total is None:
        return 0
    if isinstance(page_limit_total, (int, float)):
        return int(page_limit_total * WORDS_PER_PAGE_TNR_10PT_SINGLE)
    text = str(page_limit_total)
    match = re.search(r"\b(\d{1,3})\b", text)
    if not match:
        return 0
    return int(match.group(1)) * WORDS_PER_PAGE_TNR_10PT_SINGLE


# ============================================================
# PowerPoint extraction (shared by pp_writer + proposal_writer)
# ============================================================

def extract_text_from_pptx(file_path: str) -> str:
    """
    Pulls all readable text from a PowerPoint file: slide titles, text boxes,
    tables, and speaker notes. Labels each slide so Claude can follow the
    structure of the deck.
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("Error: python-pptx is required. Run: pip install python-pptx")
        sys.exit(1)

    prs = Presentation(file_path)
    slides_text = []

    for i, slide in enumerate(prs.slides):
        slide_content = []
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_content.append(f"[Slide {i + 1}]: {slide.shapes.title.text.strip()}")
        else:
            slide_content.append(f"[Slide {i + 1}]")

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_content.append(" | ".join(cells))
            elif hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_content.append(f"[Notes]: {notes}")

        if len(slide_content) > 1:
            slides_text.append("\n".join(slide_content))

    return "\n\n".join(slides_text)


def read_past_performance_decks(folder_path: str) -> list:
    """
    Reads all .pptx files in a past-performance folder. Returns a list of
    dicts: [{"filename": ..., "text": ...}].
    """
    pptx_files = sorted(glob.glob(os.path.join(folder_path, "*.pptx")))
    if not pptx_files:
        print(f"Error: No .pptx files found in {folder_path}")
        sys.exit(1)

    decks = []
    print(f"Found {len(pptx_files)} past performance deck(s):")
    for f in pptx_files:
        print(f"  - {os.path.basename(f)}")
        text = extract_text_from_pptx(f)
        decks.append({"filename": os.path.basename(f), "text": text})
    return decks


def format_decks_for_prompt(decks: list) -> str:
    """Format the PP deck contents into a single labeled string for the LLM."""
    combined = ""
    for deck in decks:
        combined += f"\n{'='*60}\n"
        combined += f"PAST PERFORMANCE DECK: {deck['filename']}\n"
        combined += f"{'='*60}\n\n"
        combined += deck["text"]
        combined += "\n\n"
    return combined


# ============================================================
# Markdown -> DOCX rendering
# ============================================================

def markdown_to_docx(draft_text: str, title: str, subtitle_lines: list, output_path: str) -> None:
    """
    Converts Claude's markdown draft output to a formatted Word document.
    Handles # headings, ## headings, **bold**, bullet lines, plain paragraphs,
    and horizontal-rule dividers (rendered as a tight Word rule paragraph,
    not 60 underscore characters).

    Spacing is deliberately tight (Section L typically requires 1" margins
    and every page counts):
      - 1.0" margins on all sides
      - Normal style has near-zero space-after to avoid double-gaps
      - Blank markdown lines are skipped (Word renders paragraph breaks already)
    """
    try:
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
    except ImportError:
        print("Error: python-docx is required. Run: pip install python-docx")
        sys.exit(1)

    doc = Document()

    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)

    # Tighten default paragraph spacing
    normal = doc.styles["Normal"]
    normal.paragraph_format.space_before = Pt(0)
    normal.paragraph_format.space_after = Pt(2)

    # Cover block
    h = doc.add_heading(title, level=1)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for line in subtitle_lines:
        p = doc.add_paragraph(line)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()

    def _add_horizontal_rule(doc):
        """Render an actual thin Word horizontal-rule paragraph instead of 60 underscores."""
        p = doc.add_paragraph()
        p_pr = p._p.get_or_add_pPr()
        p_bdr = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "6")
        bottom.set(qn("w:space"), "1")
        bottom.set(qn("w:color"), "BFBFBF")
        p_bdr.append(bottom)
        p_pr.append(p_bdr)
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after = Pt(0)

    for raw_line in draft_text.splitlines():
        line = raw_line.rstrip()

        if line.startswith("### "):
            doc.add_heading(line[4:].strip(), level=3)
        elif line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:].strip(), level=1)
        elif line.startswith("- ") or line.startswith("* "):
            p = doc.add_paragraph(style="List Bullet")
            _add_inline_bold(p, line[2:])
        elif line.startswith("---"):
            # Render as a real Word horizontal rule paragraph (tight; no extra space)
            _add_horizontal_rule(doc)
        elif line.strip() == "":
            # Skip blank markdown lines - Word renders paragraph breaks via spacing
            continue
        else:
            p = doc.add_paragraph()
            _add_inline_bold(p, line)

    doc.save(output_path)


def _add_inline_bold(paragraph, text: str) -> None:
    """Renders **bold** spans within a paragraph run."""
    parts = text.split("**")
    for i, part in enumerate(parts):
        run = paragraph.add_run(part)
        run.bold = (i % 2 == 1)
